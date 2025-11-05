# SPDX-License-Identifier: MIT
# Copyright (c) 2025 Elena Viter
# chat/sdk/tools/pptx_renderer.py

from __future__ import annotations
import pathlib
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any
import json
import re
from html.parser import HTMLParser
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from kdcube_ai_app.apps.chat.sdk.runtime.workdir_discovery import resolve_output_dir
import kdcube_ai_app.apps.chat.sdk.tools.md_utils as md_utils


# ============================================================================
# CONFIGURATION
# ============================================================================
EMU_PER_INCH = 914400
SLIDE_HEIGHT = Inches(7.5)
SLIDE_WIDTH = Inches(10)
MARGIN = Inches(0.5)
MAX_CONTENT_HEIGHT = Inches(6.0)  # Leave room for title

def page_width_in() -> float:
    return SLIDE_WIDTH.inches - 2 * MARGIN.inches

# Default colors
DEFAULT_COLORS = {
    'text': RGBColor(51, 51, 51),
    'primary': RGBColor(0, 102, 204),
    'subtitle': RGBColor(102, 102, 102),
}


# ============================================================================
# CSS PARSER
# ============================================================================

@dataclass
class StyleInfo:
    color: Optional[RGBColor] = None
    background: Optional[RGBColor] = None
    font_size: Optional[Pt] = None
    line_height: Optional[float] = None  # multiplier, e.g. 1.6
    padding_left: Optional[Inches] = None
    padding_top: Optional[Inches] = None
    padding_right: Optional[Inches] = None
    padding_bottom: Optional[Inches] = None
    border_color: Optional[RGBColor] = None
    border_width: Optional[Pt] = None
    border_left_color: Optional[RGBColor] = None
    border_left_width: Optional[Pt] = None
    bold: bool = False
    italic: bool = False


class CSSParser:
    def __init__(self):
        self.styles: Dict[str, Dict[str, str]] = {}

    def get_style(self, class_name: str) -> StyleInfo:
        """Get computed style for a class."""
        rules = self.styles.get(class_name, {})
        style = StyleInfo()

        if 'color' in rules:
            style.color = self._parse_color(rules['color'])

        if 'background' in rules:
            style.background = self._parse_color(rules['background'])

        if 'font-size' in rules:
            style.font_size = self._parse_font_size(rules['font-size'])

        # Parse border-bottom (for title underline)
        if 'border-bottom' in rules:
            parts = rules['border-bottom'].split()
            if len(parts) >= 3:
                try:
                    # e.g., "4px solid #0066cc"
                    width_str = parts[0].replace('px', '').replace('pt', '')
                    style.border_width = Pt(float(width_str))
                    style.border_color = self._parse_color(parts[2])
                except:
                    pass

        # Parse border-left (for callout boxes)
        if 'border-left' in rules:
            parts = rules['border-left'].split()
            if len(parts) >= 3:
                try:
                    width_str = parts[0].replace('px', '').replace('pt', '')
                    style.border_left_width = Pt(float(width_str))
                    style.border_left_color = self._parse_color(parts[2])
                except:
                    pass

        return style

    def parse(self, css_text: str):
        pattern = r'([a-zA-Z0-9_\-., #]+)\s*\{([^}]+)\}'
        for match in re.finditer(pattern, css_text):
            selector = match.group(1).strip()
            # map ".highlight" -> "highlight", "h1" -> "h1", "strong" -> "strong"
            selector = selector.replace('.', '').replace(' ', '')
            rules = {}
            for rule in match.group(2).split(';'):
                rule = rule.strip()
                if ':' in rule:
                    prop, value = rule.split(':', 1)
                    rules[prop.strip().lower()] = value.strip()
            self.styles[selector] = rules

    # NEW: parse inline style="..."
    def parse_inline_rules(self, style_attr: str) -> Dict[str, str]:
        rules = {}
        for piece in style_attr.split(';'):
            piece = piece.strip()
            if ':' in piece:
                k, v = piece.split(':', 1)
                rules[k.strip().lower()] = v.strip()
        return rules

    def _parse_color(self, color_str: str) -> Optional[RGBColor]:
        s = color_str.strip().lower()
        if s.startswith('#'):
            s = s[1:]
        if len(s) == 6:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
        if len(s) == 3:
            return RGBColor(int(s[0]*2, 16), int(s[1]*2, 16), int(s[2]*2, 16))
        return None  # do not hardcode

    def _parse_font_size(self, size_str: str) -> Optional[Pt]:
        s = size_str.strip().lower()
        if s.endswith('em'):
            return Pt(float(s[:-2]) * 16 * 0.75)  # 1em=16px; 1px=0.75pt
        if s.endswith('px'):
            return Pt(float(s[:-2]) * 0.75)
        if s.endswith('pt'):
            return Pt(float(s[:-2]))
        return None

    def _parse_length_to_inches(self, s: str) -> Optional[Inches]:
        s = s.strip().lower()
        try:
            if s.endswith('px'):  # assume 96dpi for CSS pixel
                return Inches(float(s[:-2]) / 96.0)
            if s.endswith('pt'):
                return Inches(float(s[:-2]) / 72.0)
            if s.endswith('in'):
                return Inches(float(s[:-2]))
        except:
            pass
        return None

    def _inflate_box_shorthand(self, v: str) -> Tuple[Optional[str], Optional[str], Optional[str], Optional[str]]:
        # CSS padding shorthand: t r b l
        parts = v.split()
        if not parts: return (None, None, None, None)
        if len(parts) == 1: parts = parts * 4
        if len(parts) == 2: parts = [parts[0], parts[1], parts[0], parts[1]]
        if len(parts) == 3: parts = [parts[0], parts[1], parts[2], parts[1]]
        return tuple(parts[:4])

    # NEW: merge rules from tag + classes + inline + body with precedence inline>class>tag>body
    def compute_style(self, tag: str, classes: List[str], inline_style: Optional[str]) -> StyleInfo:
        order = []
        if 'body' in self.styles:
            order.append(self.styles['body'])
        if tag in self.styles:
            order.append(self.styles[tag])
        for cls in classes:
            if cls in self.styles:
                order.append(self.styles[cls])
        if inline_style:
            order.append(self.parse_inline_rules(inline_style))

        style = StyleInfo()
        for rules in order:
            if 'color' in rules:
                style.color = self._parse_color(rules['color'])
            if 'background' in rules:
                style.background = self._parse_color(rules['background'])
            if 'background-color' in rules:
                style.background = self._parse_color(rules['background-color'])
            if 'font-size' in rules:
                style.font_size = self._parse_font_size(rules['font-size'])
            if 'line-height' in rules:
                try:
                    style.line_height = float(re.sub('[^0-9.]', '', rules['line-height']))
                except:
                    pass
            if 'padding' in rules:
                t, r, b, l = self._inflate_box_shorthand(rules['padding'])
                style.padding_top = self._parse_length_to_inches(t) if t else style.padding_top
                style.padding_right = self._parse_length_to_inches(r) if r else style.padding_right
                style.padding_bottom = self._parse_length_to_inches(b) if b else style.padding_bottom
                style.padding_left = self._parse_length_to_inches(l) if l else style.padding_left
            for side in ('padding-left','padding-right','padding-top','padding-bottom'):
                if side in rules:
                    val = self._parse_length_to_inches(rules[side])
                    if side == 'padding-left': style.padding_left = val
                    if side == 'padding-right': style.padding_right = val
                    if side == 'padding-top': style.padding_top = val
                    if side == 'padding-bottom': style.padding_bottom = val
            if 'border-bottom' in rules:
                parts = rules['border-bottom'].split()
                if len(parts) >= 3:
                    try:
                        style.border_width = Pt(float(re.sub('[^0-9.]', '', parts[0])))
                        style.border_color = self._parse_color(parts[-1])
                    except:
                        pass
            if 'border-left' in rules:
                parts = rules['border-left'].split()
                if len(parts) >= 3:
                    try:
                        style.border_left_width = Pt(float(re.sub('[^0-9.]', '', parts[0])))
                        style.border_left_color = self._parse_color(parts[-1])
                    except:
                        pass
        return style

# ============================================================================
# HTML PARSER
# ============================================================================

class StyledHTMLParser(HTMLParser):
    def __init__(self, css_parser: CSSParser):
        super().__init__()
        self.css = css_parser
        self.body_style = self.css.compute_style('body', [], None)
        self.slides = []
        self.current_section = None
        self.current_element = None
        self.current_list = None
        self.current_list_item = None
        self.in_callout = False  # NEW

        self.format_stack: List[Dict[str, Any]] = []
        self.class_stack: List[List[str]] = []

    def _current_format(self) -> Dict[str, Any]:
        return self.format_stack[-1] if self.format_stack else {'bold': False, 'italic': False, 'classes': []}

    def _push_format(self, **kwargs):
        fmt = self._current_format().copy()
        fmt.update(kwargs)
        self.format_stack.append(fmt)

    def _pop_format(self):
        if self.format_stack:
            self.format_stack.pop()
        if self.class_stack:
            self.class_stack.pop()

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        classes = attrs_dict.get('class', '').split()
        inline = attrs_dict.get('style')

        self.class_stack.append(classes)

        if tag in ('strong', 'b'):
            self._push_format(bold=True)
        elif tag in ('em', 'i'):
            self._push_format(italic=True)
        elif tag == 'span':
            self._push_format(classes=classes)

        if tag == 'section':
            self.current_section = {
                'id': attrs_dict.get('id', ''),
                'title': '',
                'subtitle': '',
                'elements': [],
                'title_style': None,
                'subtitle_style': None
            }

        elif tag == 'h1' and self.current_section is not None:
            self.current_element = {
                'type': 'h1',
                'text': '',
                'runs': [],
                'classes': classes,
                'style': self.css.compute_style('h1', classes, inline)
            }

        elif tag == 'h2' and self.current_section is not None:
            self.current_element = {
                'type': 'h2',
                'text': '',
                'runs': [],
                'classes': classes,
                'style': self.css.compute_style('h2', classes, inline)
            }

        elif tag == 'h3' and self.current_section is not None:
            self.current_element = {
                'type': 'h3',
                'text': '',
                'runs': [],
                'classes': classes,
                'style': self.css.compute_style('h3', classes, inline)
            }

        elif tag == 'p' and self.current_section is not None:
            # If we are inside a callout, keep writing into that callout element, do not create a new paragraph element
            if self.in_callout and self.current_element and self.current_element.get('type') == 'callout':
                return
            elem_style = self.css.compute_style('p', classes, inline)
            self.current_element = {
                'type': 'subtitle' if 'subtitle' in classes else 'paragraph',
                'text': '',
                'runs': [],
                'classes': classes,
                'style': elem_style
            }

        elif tag in ('ul', 'ol') and self.current_section is not None:
            self.current_list = {'type': 'list', 'items': []}

        elif tag == 'li' and self.current_list is not None:
            self.current_list_item = {'text': '', 'runs': [], 'style': self.css.compute_style('li', classes, inline)}

        elif tag == 'div' and self.current_section is not None:
            elem_style = self.css.compute_style('div', classes, inline)
            if 'highlight' in classes or elem_style.background or elem_style.border_left_color:
                self.current_element = {
                    'type': 'callout',
                    'text': '',
                    'runs': [],
                    'classes': classes,
                    'style': elem_style
                }
                self.in_callout = True

    def handle_endtag(self, tag):
        if tag in ('strong', 'b', 'em', 'i', 'span'):
            self._pop_format()

        if tag == 'section' and self.current_section is not None:
            self.slides.append(self.current_section); self.current_section = None

        elif tag == 'h1' and self.current_element:
            self.current_section['title'] = self.current_element['text'].strip()
            self.current_section['title_style'] = self.current_element['style']
            self.current_element = None

        elif tag in ('h2', 'h3') and self.current_element:
            self.current_section['elements'].append(self.current_element)
            self.current_element = None

        elif tag == 'p' and self.current_element:
            if self.in_callout and self.current_element.get('type') == 'callout':
                # nothing to close; we kept writing runs into callout
                return
            if self.current_element['type'] == 'subtitle':
                self.current_section['subtitle'] = self.current_element['text'].strip()
                self.current_section['subtitle_style'] = self.current_element['style']
            else:
                self.current_section['elements'].append(self.current_element)
            self.current_element = None

        elif tag in ('ul', 'ol') and self.current_list:
            self.current_section['elements'].append(self.current_list); self.current_list = None

        elif tag == 'li' and self.current_list_item:
            self.current_list['items'].append(self.current_list_item); self.current_list_item = None

        elif tag == 'div' and self.current_element:
            if self.current_element.get('type') == 'callout':
                self.current_section['elements'].append(self.current_element)
            self.current_element = None
            self.in_callout = False

    def handle_data(self, data):
        if not data.strip():
            return
        fmt = self._current_format()
        # Inherit color: inline class run color > element style color > body color
        color = None
        for cls in fmt.get('classes', []):
            s = self.css.get_style(cls)  # still okay for span classes
            if s.color:
                color = s.color; break
        if color is None and self.current_element:
            elem_style = self.current_element.get('style')
            if elem_style and elem_style.color:
                color = elem_style.color
        if color is None and self.body_style.color:
            color = self.body_style.color

        run = {'text': data, 'bold': fmt.get('bold', False), 'italic': fmt.get('italic', False), 'color': color}

        if self.current_list_item is not None:
            self.current_list_item['text'] += data
            self.current_list_item['runs'].append(run)
        elif self.current_element is not None:
            self.current_element['text'] += data
            if 'runs' in self.current_element:
                self.current_element['runs'].append(run)


# ============================================================================
# RENDERER WITH OVERFLOW PROTECTION
# ============================================================================
def _avg_char_width_pt(font_pt: float, bold: bool) -> float:
    # rough average, works well enough for layout
    base = font_pt * 0.52
    return base * (1.08 if bold else 1.0)

def _estimate_runs_height_in_inches(runs: List[Dict[str,Any]], font_pt: float, width_in_inches: float, line_height_mult: float = 1.2) -> Inches:
    if not runs:
        return Inches((font_pt * line_height_mult) / 72.0)
    width_pts = width_in_inches * 72.0
    text = ''.join(r['text'] for r in runs)
    # estimate chars per line using weighted average width across runs
    if not text.strip():
        return Inches((font_pt * line_height_mult) / 72.0)

    # conservative: use the smallest chars-per-line across runs
    cpl_candidates = []
    for r in runs:
        fw = _avg_char_width_pt(font_pt, r.get('bold', False))
        if fw <= 0: continue
        cpl_candidates.append(int(max(1, width_pts / fw)))
    cpl = min(cpl_candidates) if cpl_candidates else int(max(1, width_pts / _avg_char_width_pt(font_pt, False)))

    # rough wrap count
    lines = 0
    for paragraph in text.split('\n'):
        length = max(1, len(paragraph))
        lines += int((length + cpl - 1) // cpl)
    if lines < 1: lines = 1
    height_pts = lines * (font_pt * line_height_mult)
    return Inches(height_pts / 72.0)


def estimate_content_height(elements: List[Dict], base_font_size: Pt) -> float:
    total_in = 0.0
    page_w_in = SLIDE_WIDTH.inches - 2 * MARGIN.inches  # <-- float inches

    for elem in elements:
        t = elem.get('type')
        style: StyleInfo = elem.get('style') or StyleInfo()

        if t in ('h2', 'h3'):
            fs = (style.font_size or (Pt(28) if t == 'h2' else Pt(22))).pt
            h = _estimate_runs_height_in_inches(
                elem.get('runs', []),
                fs,
                page_w_in,
                (style.line_height or 1.2),
            )
            total_in += h.inches + 0.1

        elif t == 'list':
            fs = base_font_size.pt
            bullets_w_in = page_w_in - 0.25  # 0.25" indent on the left
            for it in elem.get('items', []):
                runs = it.get('runs', [])
                h = _estimate_runs_height_in_inches(runs, fs, bullets_w_in, 1.3)
                total_in += h.inches + 0.06
            total_in += 0.1

        elif t == 'callout':
            fs = (style.font_size or Pt(16)).pt
            pad_l_in = (style.padding_left or Inches(0.2)).inches
            pad_r_in = (style.padding_right or Inches(0.1)).inches
            pad_t_in = (style.padding_top or Inches(0.1)).inches
            pad_b_in = (style.padding_bottom or Inches(0.1)).inches

            content_w_in = page_w_in - pad_l_in - pad_r_in
            h = _estimate_runs_height_in_inches(
                elem.get('runs', []),
                fs,
                content_w_in,
                (style.line_height or 1.3),
            )
            total_in += max(h.inches + pad_t_in + pad_b_in, 0.6) + 0.15

    return total_in


def render_slide(prs: Presentation, slide_data: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text = slide_data.get('title', '')
    title_style: StyleInfo = slide_data.get('title_style') or StyleInfo()
    subtitle_text = slide_data.get('subtitle', '')
    subtitle_style: StyleInfo = slide_data.get('subtitle_style') or StyleInfo()
    elements = slide_data.get('elements', [])

    page_w_in = SLIDE_WIDTH.inches - 2 * MARGIN.inches  # float inches

    estimated_height = estimate_content_height(elements, Pt(18))
    scale_factor = 1.0
    if estimated_height > MAX_CONTENT_HEIGHT.inches:
        scale_factor = max(0.7, MAX_CONTENT_HEIGHT.inches / estimated_height)

    y = MARGIN

    # TITLE
    title_h = Inches(0.8 * scale_factor)
    title_box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - MARGIN * 2, title_h)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title_text
    r.font.size = title_style.font_size or Pt(36 * scale_factor)
    r.font.bold = True
    if title_style.color:
        r.font.color.rgb = title_style.color
    y += title_h

    # underline from CSS
    if title_style.border_color and title_style.border_width:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN, y,
            SLIDE_WIDTH - MARGIN * 2, title_style.border_width
        )
        line.fill.solid()
        line.fill.fore_color.rgb = title_style.border_color
        line.line.fill.background()
        y += Inches(0.12 * scale_factor)

    # SUBTITLE
    if subtitle_text:
        sub_h = Inches(0.45 * scale_factor)
        subtitle_box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - MARGIN * 2, sub_h)
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = subtitle_text
        r.font.size = subtitle_style.font_size or Pt(18 * scale_factor)
        r.font.italic = subtitle_style.italic
        if subtitle_style.color:
            r.font.color.rgb = subtitle_style.color
        y += sub_h + Inches(0.12 * scale_factor)

    # CONTENT
    base_font_pt = 18 * scale_factor

    for elem in elements:
        elem_type = elem.get('type')
        style: StyleInfo = elem.get('style') or StyleInfo()
        runs = elem.get('runs', [])

        if elem_type in ('h2', 'h3'):
            fs = (style.font_size.pt if style.font_size else (28 if elem_type == 'h2' else 22)) * scale_factor
            box_h = _estimate_runs_height_in_inches(runs, fs, page_w_in, (style.line_height or 1.2))
            # drawing width stays as Length
            box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - MARGIN * 2, box_h)
            tf = box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            for run_data in runs:
                rr = p.add_run(); rr.text = run_data['text']; rr.font.size = Pt(fs); rr.font.bold = True
                if run_data.get('color'):
                    rr.font.color.rgb = run_data['color']
                elif style.color:
                    rr.font.color.rgb = style.color
            y += box_h + Inches(0.08 * scale_factor)

        elif elem_type == 'list':
            # keep separate measurement width vs drawing width
            w_len = SLIDE_WIDTH - MARGIN * 2 - Inches(0.25)     # Length/EMU for drawing
            w_in  = page_w_in - 0.25                            # float inches for measurement
            for item in elem.get('items', []):
                runs_i = item.get('runs', [])
                fs = base_font_pt
                lh = 1.3
                h = _estimate_runs_height_in_inches(runs_i, fs, w_in, lh)
                box = slide.shapes.add_textbox(MARGIN + Inches(0.25), y, w_len, h)
                tf = box.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.level = 0
                for run_data in runs_i:
                    rr = p.add_run(); rr.text = run_data['text']; rr.font.size = Pt(fs)
                    rr.font.bold = run_data.get('bold', False)
                    if run_data.get('color'):
                        rr.font.color.rgb = run_data['color']
                y += h + Inches(0.06 * scale_factor)
            y += Inches(0.06 * scale_factor)

        elif elem_type == 'callout':
            fs = (style.font_size.pt if style.font_size else 16) * scale_factor
            pad_l = style.padding_left or Inches(0.2)
            pad_r = style.padding_right or Inches(0.1)
            pad_t = style.padding_top or Inches(0.1)
            pad_b = style.padding_bottom or Inches(0.1)

            content_w_in = page_w_in - pad_l.inches - pad_r.inches
            content_w_len = SLIDE_WIDTH - MARGIN * 2 - pad_l - pad_r  # for drawing

            content_h = _estimate_runs_height_in_inches(runs, fs, content_w_in, (style.line_height or 1.3))
            box_h_in = max(content_h.inches + pad_t.inches + pad_b.inches, 0.6)
            box_h = Inches(box_h_in)

            if style.background:
                bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, SLIDE_WIDTH - MARGIN * 2, box_h)
                bg.fill.solid(); bg.fill.fore_color.rgb = style.background
                bg.line.fill.background()

            if style.border_left_color and style.border_left_width:
                border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, style.border_left_width, box_h)
                border.fill.solid(); border.fill.fore_color.rgb = style.border_left_color
                border.line.fill.background()

            text_box = slide.shapes.add_textbox(MARGIN + pad_l, y + pad_t, content_w_len, box_h - pad_t - pad_b)
            tf = text_box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            for run_data in runs:
                rr = p.add_run(); rr.text = run_data['text']; rr.font.size = Pt(fs); rr.font.bold = run_data.get('bold', False)
                if run_data.get('color'):
                    rr.font.color.rgb = run_data['color']
                elif style.color:
                    rr.font.color.rgb = style.color

            y += box_h + Inches(0.12 * scale_factor)

# ============================================================================
# MAIN
# ============================================================================

def render_pptx(
        path: str,
        content_md: str = "",
        content_html: str = "",
        *,
        title: Optional[str] = None,
        base_dir: Optional[str] = None,
        sources: Optional[str] = None,
        resolve_citations: bool = False,
        include_sources_slide: bool = False
) -> str:
    """Render PowerPoint from HTML."""

    outdir = resolve_output_dir()
    filename = Path(path).name
    if not filename.endswith('.pptx'):
        filename += '.pptx'
    outfile = outdir / filename
    outfile.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    if content_html and content_html.strip():
        # Parse CSS
        css_parser = CSSParser()
        style_match = re.search(r'<style[^>]*>(.*?)</style>', content_html, re.DOTALL | re.I)
        if style_match:
            css_parser.parse(style_match.group(1))

            # Debug: print parsed styles
            print("=== PARSED CSS ===")
            for selector, rules in css_parser.styles.items():
                print(f"{selector}: {rules}")

        # Parse HTML
        html_parser = StyledHTMLParser(css_parser)
        html_parser.feed(content_html)

        slides_data = html_parser.slides

        print(f"=== PARSED {len(slides_data)} SLIDES ===")

        # Render ALL slides
        for slide_data in slides_data:
            print(f"Rendering slide: {slide_data.get('title')}")
            render_slide(prs, slide_data)

    prs.save(str(outfile))
    return filename