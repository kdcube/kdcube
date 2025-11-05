# SPDX-License-Identifier: MIT
# Copyright (c) 2025 Elena Viter

# chat/sdk/tools/pptx_renderer.py
# FINAL FIXED VERSION - H3 headers, proper height tracking, overflow prevention

from __future__ import annotations

import pathlib
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any
import json
import re
from html.parser import HTMLParser
from typing import Dict, List, Tuple, Optional, Any
import tinycss2

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

from kdcube_ai_app.apps.chat.sdk.runtime.workdir_discovery import resolve_output_dir
import kdcube_ai_app.apps.chat.sdk.tools.md_utils as md_utils


# Layout heuristics (tuned to avoid overlap)
TITLE_BAND_IN   = 0.90   # reserved height below slide title
BOTTOM_MARGIN_IN = 0.50  # keep off the footer area
BLOCK_GAP_IN     = 0.12  # vertical rhythm between blocks

# line-height multipliers (approx)
LINE_SPACING = 1.25
EMU_PER_IN = 914400

def _emu_to_in(v: int) -> float:
    return float(v) / EMU_PER_IN

def _in_to_emu(inches: float) -> int:
    return int(round(inches * EMU_PER_IN))

def _chars_per_line(width_in: float, font_pt: float, indent_level: int = 0) -> int:
    """
    Very conservative char-per-line estimate.
    Rough model: avg char width ≈ 0.55em; with bullets/indent we reduce effective width.
    """
    # effective width reduction per indent level (~0.35" per level)
    indent_in = max(0, indent_level) * 0.35
    eff_in = max(1.0, width_in - indent_in)
    # points per inch = 72; approx chars per em = 2; fudge factor 0.8 for safety
    chars = int((eff_in * 72 / font_pt) * 2 * 0.8)
    # cap to a sane range
    return max(25, min(95, chars))

def _lines_for_text(text: str, width_in: float, font_pt: float, indent_level: int = 0) -> int:
    # strip markdown emphasis markers for width calc
    stripped = re.sub(r"(\*\*|\*)", "", text)
    cpl = _chars_per_line(width_in, font_pt, indent_level)
    # very conservative break: count words/characters
    return max(1, (len(stripped) + cpl - 1) // cpl)

def _pt_to_in(pts: float) -> float:
    return pts / 72.0

PARA_SPACE_AFTER_PT = 2      # you already use this in _style_paragraph
TF_MARGIN_IN = 0.10          # 0.05 top + 0.05 bottom (from _add_textbox)
HEIGHT_FUDGE = 1.10          # safety multiplier for long tokens/URLs

def _estimate_text_block_height(lines: List[str], width_in: float, font_pt: float) -> float:
    """
    Estimate height (in inches) for a list of paragraph lines (each becomes its own paragraph).
    Accounts for wrapping, per-paragraph spacing, text frame margins, and a safety fudge.
    """
    if not lines:
        return 0.0

    line_height_in = _pt_to_in(font_pt) * LINE_SPACING
    para_space_in = _pt_to_in(PARA_SPACE_AFTER_PT)

    total_visual_lines = 0
    nonempty_paras = 0

    for ln in lines:
        if not ln.strip():
            continue
        lvl, txt = _parse_bullet_level(ln)
        # be more conservative on width for deep indents
        visual_lines = _lines_for_text(txt, width_in * 0.95, font_pt, lvl)
        total_visual_lines += visual_lines
        nonempty_paras += 1

    if nonempty_paras == 0:
        return 0.0

    content_h = total_visual_lines * line_height_in
    # add paragraph spacing between paragraphs (n-1 gaps)
    spacing_h = max(0, nonempty_paras - 1) * para_space_in
    # add textframe top+bottom margins
    margins_h = TF_MARGIN_IN

    return (content_h + spacing_h + margins_h) * HEIGHT_FUDGE
# GPT


_LINK_RE = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")
_CIT_RE  = re.compile(r"\[\[S:(\d+)\]\]")

def _outdir() -> pathlib.Path:
    return resolve_output_dir()

def _basename_only(path: str, default_ext: str = ".pptx") -> str:
    name = Path(path).name
    if default_ext and not name.lower().endswith(default_ext):
        name += default_ext
    return name

def _ensure_parent(p: Path) -> None:
    p.parent.mkdir(parents=True, exist_ok=True)

def _domain_of(url: str) -> str:
    try:
        from urllib.parse import urlparse
        net = urlparse(url).netloc
        return net or url
    except Exception:
        return url

def _split_markdown_sections(md: str) -> List[Tuple[str, List[str]]]:
    """Split markdown into sections by ## headings."""
    lines = (md or "").splitlines()
    slides: List[Tuple[str, List[str]]] = []
    cur_title: Optional[str] = None
    cur_body: List[str] = []

    for ln in lines:
        if ln.startswith("## "):
            if cur_title is not None:
                slides.append((cur_title.strip(), cur_body))
            cur_title = ln[3:]
            cur_body = []
        elif ln.startswith("# "):
            if cur_title is None and not slides:
                cur_title = ln[2:]
                cur_body = []
            else:
                cur_body.append(ln)
        else:
            cur_body.append(ln)

    if cur_title is None:
        nonempty = next((l for l in lines if l.strip()), "Slides")
        cur_title = nonempty.lstrip("# ").strip() or "Slides"

    slides.append((cur_title.strip(), cur_body))
    return slides

# Styling
PALETTE = {
    "fg": RGBColor(20, 24, 31),
    "muted": RGBColor(95, 106, 121),
    "accent": RGBColor(31, 111, 235),
    "quote_bg": RGBColor(245, 247, 250),
    "code_bg": RGBColor(250, 250, 252),
    "rule": RGBColor(220, 224, 230),
    # "table_row_bg": RGBColor(255, 255, 255),      # even rows (base)
    "table_header_bg": RGBColor(240, 244, 252),
    # "table_header_fg": RGBColor(20, 24, 31),   # dark text in header
    # "table_row_alt_bg": RGBColor(232, 237, 245)  # light zebra stripe
}

TYPE_SCALE = {
    "title": Pt(36),
    "slide_title": Pt(26),
    "h3": Pt(20),       # For ### headers
    "body": Pt(15),     # Slightly smaller for better fit
    "code": Pt(12),
    "caption": Pt(11),
}

PAGE = {
    "content_left": Inches(0.8),
    "content_top": Inches(1.4),
    "content_width": Inches(8.4),  # 10.0 - 0.8 - 0.8
    "slide_height": Inches(7.5),
}

MONO_FALLBACK = "Consolas"

# Regex patterns
_CODE_FENCE_RE = re.compile(r"^```(\w+)?\s*$")
_TABLE_ROW_RE = re.compile(r"^\s*\|.+\|\s*$")
_BLOCKQUOTE_RE = re.compile(r"^\s*>\s+(.*)$")

def _style_run(run, *, size: Pt, bold=False, italic=False, color: RGBColor = None, mono=False):
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    if mono:
        run.font.name = MONO_FALLBACK
    if color:
        run.font.color.rgb = color

def _style_paragraph(p, *, level=0, space_after=Pt(6), align=PP_ALIGN.LEFT):
    p.level = max(0, min(level, 4))
    p.space_after = space_after
    p.line_spacing = 1.15
    p.alignment = align

def _add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tb

def _add_title_slide(prs: Presentation, text: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tbox = _add_textbox(slide, PAGE["content_left"], Inches(2.0), PAGE["content_width"], Inches(2.0))
    p = tbox.text_frame.paragraphs[0]
    _style_paragraph(p, align=PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = (text or "Presentation").strip()
    _style_run(r, size=TYPE_SCALE["title"], bold=True, color=PALETTE["fg"])

def _add_slide_title(slide, title: str):
    tbox = _add_textbox(slide, PAGE["content_left"], PAGE["content_top"], PAGE["content_width"], Inches(0.6))
    p = tbox.text_frame.paragraphs[0]
    _style_paragraph(p, align=PP_ALIGN.LEFT, space_after=Pt(2))
    r = p.add_run()
    r.text = (title or "").strip()
    _style_run(r, size=TYPE_SCALE["slide_title"], bold=True, color=PALETTE["fg"])

def _parse_bullet_level(line: str) -> tuple[int, str]:
    """Extract bullet level and text."""
    m = re.match(r"^(\s*)([-*]|\d+\.)\s+(.*)$", line)
    if not m:
        return 0, line.strip()
    spaces, _bullet, text = m.groups()
    lvl = min(len(spaces) // 2, 4)
    return lvl, text.strip()

def _is_table_separator(cell: str) -> bool:
    """Check if a table cell is a separator (contains dashes)."""
    stripped = cell.strip()
    if not stripped:
        return False
    dash_count = stripped.count('-')
    return dash_count >= 3 and dash_count >= len(stripped) - 2

def _parse_table(lines: List[str]) -> Optional[List[List[str]]]:
    """Parse markdown table into rows/columns."""
    table_lines = [ln.strip() for ln in lines if _TABLE_ROW_RE.match(ln)]
    if len(table_lines) < 2:
        return None

    def split_row(line: str) -> List[str]:
        return [cell.strip() for cell in line.strip('|').split('|')]

    rows = [split_row(line) for line in table_lines]

    if len(rows) < 2:
        return None

    # Check if second row is separator
    separator_row = rows[1]
    if not all(_is_table_separator(cell) for cell in separator_row):
        return None

    # Return header + data rows (skip separator)
    header = rows[0]
    data_rows = rows[2:] if len(rows) > 2 else []

    return [header] + data_rows

def _emit_text_with_formatting(paragraph, text: str, sources_map: Dict[int, Dict[str, str]], resolve_citations: bool):
    """Add text to paragraph with bold/italic/links/citations."""
    # Split by bold first
    parts = re.split(r"(\*\*[^*]+\*\*)", text)

    for part in parts:
        if not part:
            continue

        if part.startswith("**") and part.endswith("**"):
            # Bold text
            inner = part[2:-2]
            # Check for italic within bold
            italic_parts = re.split(r"(\*[^*]+\*)", inner)
            for ipart in italic_parts:
                if ipart.startswith("*") and ipart.endswith("*"):
                    r = paragraph.add_run()
                    r.text = ipart[1:-1]
                    _style_run(r, size=TYPE_SCALE["body"], bold=True, italic=True, color=PALETTE["fg"])
                elif ipart:
                    r = paragraph.add_run()
                    r.text = ipart
                    _style_run(r, size=TYPE_SCALE["body"], bold=True, color=PALETTE["fg"])
        else:
            # Check for italic, links, citations
            italic_parts = re.split(r"(\*[^*]+\*)", part)
            for ipart in italic_parts:
                if not ipart:
                    continue

                if ipart.startswith("*") and ipart.endswith("*"):
                    r = paragraph.add_run()
                    r.text = ipart[1:-1]
                    _style_run(r, size=TYPE_SCALE["body"], italic=True, color=PALETTE["fg"])
                else:
                    # Handle links and citations
                    _emit_links_and_citations(paragraph, ipart, sources_map, resolve_citations)

def _emit_links_and_citations(paragraph, text: str, sources_map: Dict[int, Dict[str, str]], resolve_citations: bool):
    """Handle links [text](url) and citations [[S:n]]."""
    while text:
        m_link = _LINK_RE.search(text)
        m_cit = _CIT_RE.search(text) if resolve_citations else None

        matches = []
        if m_link:
            matches.append(("link", m_link))
        if m_cit:
            matches.append(("cit", m_cit))

        if not matches:
            # Plain text
            if text:
                r = paragraph.add_run()
                r.text = text
                _style_run(r, size=TYPE_SCALE["body"], color=PALETTE["fg"])
            return

        # Get earliest match
        kind, m = min(matches, key=lambda x: x[1].start())

        # Add text before match
        if m.start() > 0:
            r = paragraph.add_run()
            r.text = text[:m.start()]
            _style_run(r, size=TYPE_SCALE["body"], color=PALETTE["fg"])

        # Add match
        if kind == "link":
            r = paragraph.add_run()
            r.text = m.group(1)
            _style_run(r, size=TYPE_SCALE["body"], color=PALETTE["accent"])
            try:
                r.hyperlink.address = m.group(2)
            except:
                pass
        else:  # citation
            sid = int(m.group(1))
            rec = sources_map.get(sid, {})
            label = rec.get("title") or f"[{sid}]"
            url = rec.get("url", "")
            r = paragraph.add_run()
            r.text = label
            _style_run(r, size=TYPE_SCALE["body"], color=PALETTE["accent"])
            if url:
                try:
                    r.hyperlink.address = url
                except:
                    pass

        text = text[m.end():]

def _new_content_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, title)

    # Keep anchors as EMUs
    top_anchor = PAGE["content_top"] + Inches(TITLE_BAND_IN)          # EMU int
    bottom_anchor = PAGE["slide_height"] - Inches(BOTTOM_MARGIN_IN)   # EMU int

    # Available height *in inches* for our estimators
    available_height_in = _emu_to_in(bottom_anchor - top_anchor)

    y_offset_in = 0.0  # we track runtime offset in inches, convert to EMU only when placing
    return slide, top_anchor, available_height_in, y_offset_in

def _peek_next_block(body_lines: List[str], start_idx: int) -> Tuple[str, int, List[str]]:
    """
    Return (kind, end_index, payload_lines) starting from start_idx.
    kind in {"h3","code","table","quote","text","blank"}
    end_index is the index AFTER the block.
    """
    n = len(body_lines)
    i = start_idx
    if i >= n:
        return "blank", i, []

    ln = body_lines[i]
    if not ln.strip():
        return "blank", i+1, []

    if ln.startswith("### "):
        return "h3", i+1, [ln[4:].strip()]

    if _CODE_FENCE_RE.match(ln):
        i += 1
        code = []
        while i < n and not _CODE_FENCE_RE.match(body_lines[i]):
            code.append(body_lines[i]); i += 1
        if i < n and _CODE_FENCE_RE.match(body_lines[i]):
            i += 1
        return "code", i, code

    if _TABLE_ROW_RE.match(ln):
        table_lines = []
        while i < n and _TABLE_ROW_RE.match(body_lines[i]):
            table_lines.append(body_lines[i]); i += 1
        return "table", i, table_lines

    if _BLOCKQUOTE_RE.match(ln):
        quotes = []
        while i < n and _BLOCKQUOTE_RE.match(body_lines[i]):
            m = _BLOCKQUOTE_RE.match(body_lines[i])
            quotes.append(m.group(1)); i += 1
        return "quote", i, quotes

    # text block
    text = []
    while i < n:
        curr = body_lines[i]
        if not curr.strip():
            i += 1
            continue
        if curr.startswith("### ") or _CODE_FENCE_RE.match(curr) or _TABLE_ROW_RE.match(curr) or _BLOCKQUOTE_RE.match(curr):
            break
        text.append(curr); i += 1
    return "text", i, text

def _count_lines_height(lines: List[str], width_in: float, font_pt: float) -> Tuple[int, float]:
    """
    Return (visual_lines_count, height_in) for the given lines using same estimator pieces.
    """
    if not lines:
        return 0, 0.0
    line_height_in = _pt_to_in(font_pt) * LINE_SPACING
    para_space_in = _pt_to_in(PARA_SPACE_AFTER_PT)

    total_visual = 0
    nonempty = 0
    for ln in lines:
        if not ln.strip():
            continue
        lvl, txt = _parse_bullet_level(ln)
        total_visual += _lines_for_text(txt, width_in * 0.95, font_pt, lvl)
        nonempty += 1

    if nonempty == 0:
        return 0, 0.0
    content_h = total_visual * line_height_in
    spacing_h = max(0, nonempty - 1) * para_space_in
    margins_h = TF_MARGIN_IN
    return total_visual, (content_h + spacing_h + margins_h) * HEIGHT_FUDGE

def _pack_lines_to_height(lines: List[str], width_in: float, font_pt: float, max_h_in: float) -> int:
    """
    Return how many leading lines from `lines` can fit into `max_h_in`.
    Greedy: add one line at a time with conservative height calc.
    """
    if max_h_in <= 0:
        return 0
    lo, hi = 0, len(lines)
    # binary search for speed on big chunks
    while lo < hi:
        mid = (lo + hi + 1) // 2
        _, h = _count_lines_height(lines[:mid], width_in, font_pt)
        if h <= max_h_in:
            lo = mid
        else:
            hi = mid - 1
    return lo


def _render_section_across_slides(
        prs: Presentation,
        title: str,
        body_lines: List[str],
        sources_map: Dict[int, Dict[str, str]],
        resolve_citations: bool
) -> None:
    """
    Robust renderer that splits a section across multiple slides as needed.
    Prevents overlaps by estimating block heights conservatively and
    pushing blocks that don't fit to the next slide.
    """
    # initial slide
    slide, top_anchor, avail_in, y = _new_content_slide(prs, title)

    i = 0
    n = len(body_lines)

    def need_new_slide(cont: bool = True):
        nonlocal slide, top_anchor, avail_in, y
        cont_title = f"{title} (cont.)" if cont else title
        slide, top_anchor, avail_in, y = _new_content_slide(prs, cont_title)

    while i < n:
        # skip blank lines
        if not body_lines[i].strip():
            i += 1
            continue

        # --- H3 header ---
        if body_lines[i].startswith("### "):
            header_text = body_lines[i][4:].strip()
            header_h = 0.42  # header band (incl. spacing/padding)

            # Look ahead to the *next* block and estimate its full height.
            kind, nxt_end, payload = _peek_next_block(body_lines, i+1)
            need_after_h = 0.0
            if kind == "text":
                need_after_h = _estimate_text_block_height(payload, _emu_to_in(PAGE["content_width"]), TYPE_SCALE["body"].pt)
            elif kind == "code":
                need_after_h = 0.40 + max(1, min(20, len(payload))) * 0.18
            elif kind == "table":
                tbl = _parse_table(payload)
                if tbl:
                    rows = len(tbl)
                    need_after_h = 0.30 + rows * 0.45
            elif kind == "quote":
                qh = _estimate_text_block_height(payload, _emu_to_in(PAGE["content_width"]), TYPE_SCALE["body"].pt)
                need_after_h = max(0.60, qh + 0.20)
            else:
                need_after_h = 0.30  # minimal stub if blank/unknown

            # If header + next block won't fit, start a new slide *before* placing header
            if y + header_h + need_after_h > avail_in:
                need_new_slide(cont=True)

            # place header
            tbox = _add_textbox(slide, PAGE["content_left"], top_anchor + Inches(y),
                                PAGE["content_width"], Inches(header_h - 0.10))
            tf = tbox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            _style_paragraph(p, space_after=Pt(1))
            r = p.add_run()
            r.text = header_text
            _style_run(r, size=TYPE_SCALE["h3"], bold=True, color=PALETTE["fg"])
            y += header_h + BLOCK_GAP_IN
            i += 1
            continue

        # --- Code fence ---
        if _CODE_FENCE_RE.match(body_lines[i]):
            i += 1
            code = []
            while i < n and not _CODE_FENCE_RE.match(body_lines[i]):
                code.append(body_lines[i])
                i += 1
            # consume closing fence if present
            if i < n and _CODE_FENCE_RE.match(body_lines[i]):
                i += 1

            lines = code[:20]
            per_line = 0.18
            needed = 0.40 + max(1, len(lines)) * per_line
            if y + needed > avail_in:
                need_new_slide(cont=True)

            rect = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                PAGE["content_left"],
                top_anchor + Inches(y),
                PAGE["content_width"],
                Inches(needed - 0.06)
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = PALETTE["code_bg"]
            rect.line.color.rgb = PALETTE["rule"]
            rect.line.width = Pt(0.5)

            tf = rect.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.12)
            tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.10)
            tf.margin_bottom = Inches(0.10)

            p = tf.paragraphs[0]
            for j, code_line in enumerate(lines):
                if j > 0:
                    p = tf.add_paragraph()
                r = p.add_run()
                r.text = code_line
                _style_run(r, size=TYPE_SCALE["code"], mono=True, color=PALETTE["fg"])

            y += needed + BLOCK_GAP_IN
            continue

        # --- Table block ---
        if _TABLE_ROW_RE.match(body_lines[i]):
            table_lines = []
            while i < n and _TABLE_ROW_RE.match(body_lines[i]):
                table_lines.append(body_lines[i])
                i += 1

            table_data = _parse_table(table_lines)
            if not table_data:
                continue

            rows, cols = len(table_data), len(table_data[0])
            row_h = 0.45
            needed = 0.30 + rows * row_h
            if y + needed > avail_in:
                need_new_slide(cont=True)

            shape = slide.shapes.add_table(
                rows, cols,
                PAGE["content_left"],
                top_anchor + Inches(y),
                PAGE["content_width"],
                Inches(needed - 0.10)
            )
            tbl = shape.table

            # column widths
            col_width = int(PAGE["content_width"] / cols)
            for j in range(cols):
                tbl.columns[j].width = col_width

            # taller header row
            try:
                tbl.rows[0].height = Inches(0.55)
            except Exception:
                pass

            # fill cells + styles
            for r_idx, row_data in enumerate(table_data):
                for c_idx, cell_text in enumerate(row_data):
                    cell = tbl.cell(r_idx, c_idx)
                    cell.text = cell_text

                    # Make header background light
                    if r_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = PALETTE["table_header_bg"]

                    # text frame hygiene
                    tf = cell.text_frame
                    tf.word_wrap = True
                    tf.auto_size = MSO_AUTO_SIZE.NONE
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                    for p in tf.paragraphs:
                        p.alignment = PP_ALIGN.LEFT
                        for rr in p.runs:
                            rr.font.size = TYPE_SCALE["body"]
                            if r_idx == 0:
                                rr.font.bold = True
                                rr.font.color.rgb = PALETTE["fg"]   # dark text in header
                            else:
                                # Leave body text color/weight as-is so user/theme/zebra can control it later
                                pass

            y += needed + BLOCK_GAP_IN
            continue

        # --- Blockquote ---
        if _BLOCKQUOTE_RE.match(body_lines[i]):
            quotes = []
            while i < n and _BLOCKQUOTE_RE.match(body_lines[i]):
                m = _BLOCKQUOTE_RE.match(body_lines[i])
                quotes.append(m.group(1))
                i += 1

            # estimate quote height as normal text with italics
            q_height = _estimate_text_block_height(quotes, PAGE["content_width"].inches, TYPE_SCALE["body"].pt)
            q_height = max(0.60, q_height + 0.20)  # padding

            if y + q_height > avail_in:
                need_new_slide(cont=True)

            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                PAGE["content_left"],
                top_anchor + Inches(y),
                PAGE["content_width"],
                Inches(q_height - 0.08)
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = PALETTE["quote_bg"]
            rect.line.fill.background()

            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                PAGE["content_left"],
                top_anchor + Inches(y),
                Inches(0.08),
                Inches(q_height - 0.08)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = PALETTE["rule"]
            bar.line.fill.background()

            tf = rect.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "\n".join(quotes)
            _style_run(r, size=TYPE_SCALE["body"], italic=True, color=PALETTE["muted"])

            y += q_height + BLOCK_GAP_IN
            continue

        # --- Regular text block (collect until next special) ---
        kind, j, chunk = _peek_next_block(body_lines, i)
        if kind != "text":
            # nothing to render here
            i = max(i + 1, j)
            continue

        width_in = _emu_to_in(PAGE["content_width"])
        font_pt = TYPE_SCALE["body"].pt

        # how much space remains
        remain_in = max(0.0, avail_in - y)

        # If nothing fits on this slide, start a new one
        if remain_in < 0.35:
            need_new_slide(cont=True)
            remain_in = avail_in

        # Pack as many lines as we can into the remaining space
        fit_count = _pack_lines_to_height(chunk, width_in, font_pt, remain_in - BLOCK_GAP_IN)
        if fit_count == 0:
            need_new_slide(cont=True)
            remain_in = avail_in
            fit_count = _pack_lines_to_height(chunk, width_in, font_pt, remain_in - BLOCK_GAP_IN)

        to_render = chunk[:fit_count]
        _, est_h = _count_lines_height(to_render, width_in, font_pt)

        tbox = _add_textbox(
            slide,
            PAGE["content_left"],
            top_anchor + Inches(y),
            PAGE["content_width"],
            Inches(est_h)
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        p = tf.paragraphs[0]
        first = True
        for line in to_render:
            if not first:
                p = tf.add_paragraph()
            first = False
            lvl, txt = _parse_bullet_level(line)
            _style_paragraph(p, level=lvl, space_after=Pt(PARA_SPACE_AFTER_PT))
            _emit_text_with_formatting(p, txt, sources_map, resolve_citations)

        y += est_h + BLOCK_GAP_IN
        i += fit_count  # advance within the same logical block; we'll loop back for the remainder
        continue

        # # --- Regular text block (collect until next special) ---
        # chunk = []
        # j = i
        # while j < n:
        #     curr = body_lines[j]
        #     if not curr.strip():
        #         j += 1
        #         continue
        #     if curr.startswith("### ") or _CODE_FENCE_RE.match(curr) or _TABLE_ROW_RE.match(curr) or _BLOCKQUOTE_RE.match(curr):
        #         break
        #     chunk.append(curr)
        #     j += 1
        #
        # # nothing to render
        # if not chunk:
        #     i = max(i + 1, j)
        #     continue
        #
        # # estimate height conservatively for body text
        # est_h = _estimate_text_block_height(chunk, PAGE["content_width"].inches, TYPE_SCALE["body"].pt)
        # # if it doesn't fit, move to a new slide (but re-render same chunk)
        # if y + est_h > avail_in:
        #     need_new_slide(cont=True)
        #
        # # render the text chunk
        # tbox = _add_textbox(
        #     slide,
        #     PAGE["content_left"],
        #     top_anchor + Inches(y),
        #     PAGE["content_width"],
        #     Inches(est_h - 0.06)
        # )
        # tf = tbox.text_frame
        # tf.word_wrap = True
        # # DISABLE AUTOSIZE
        # tf.auto_size = MSO_AUTO_SIZE.NONE
        # p = tf.paragraphs[0]
        # first = True
        # for line in chunk:
        #     if not first:
        #         p = tf.add_paragraph()
        #     first = False
        #     lvl, txt = _parse_bullet_level(line)
        #     _style_paragraph(p, level=lvl, space_after=Pt(2))
        #     _emit_text_with_formatting(p, txt, sources_map, resolve_citations)
        #
        # y += est_h + BLOCK_GAP_IN
        # i = j  # advance to the next unread line

def _add_sources_slide(prs: Presentation, sources_map: Dict[int, Dict[str, str]], order: List[int]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "Sources")

    tbox = _add_textbox(
        slide,
        PAGE["content_left"],
        PAGE["content_top"] + Inches(0.8),
        PAGE["content_width"],
        Inches(5.0)
    )

    tf = tbox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]

    for idx, sid in enumerate(order):
        if idx > 0:
            p = tf.add_paragraph()

        src = sources_map.get(sid, {})
        title = src.get("title", f"Source {sid}")
        url = src.get("url", "")

        _style_paragraph(p, space_after=Pt(4))

        r1 = p.add_run()
        r1.text = f"[{sid}] "
        _style_run(r1, size=TYPE_SCALE["body"], bold=True, color=PALETTE["fg"])

        r2 = p.add_run()
        r2.text = title
        _style_run(r2, size=TYPE_SCALE["body"], color=PALETTE["fg"])

        if url:
            r3 = p.add_run()
            r3.text = f" ({_domain_of(url)})"
            _style_run(r3, size=TYPE_SCALE["caption"], color=PALETTE["accent"])


## HTML RENDERING
from html.parser import HTMLParser
class HTMLSlideParser(HTMLParser):
    """Parse HTML sections into slide structures."""

    def __init__(self):
        super().__init__()
        self.slides = []
        self.current_section = None
        self.current_element = None
        self.current_list = None
        self.current_list_item = None
        self.tag_stack = []

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        self.tag_stack.append(tag)

        if tag == 'section':
            # Start new slide
            self.current_section = {
                'id': attrs_dict.get('id', ''),
                'title': '',
                'elements': []
            }

        elif tag in ('h1', 'h2', 'h3') and self.current_section is not None:
            # Slide title
            self.current_element = {'type': 'heading', 'level': int(tag[1]), 'text': ''}

        elif tag == 'p' and self.current_section is not None:
            # Paragraph
            self.current_element = {'type': 'paragraph', 'text': ''}

        elif tag in ('ul', 'ol') and self.current_section is not None:
            # Start list
            self.current_list = {
                'type': 'list',
                'ordered': tag == 'ol',
                'items': []
            }

        elif tag == 'li' and self.current_list is not None:
            # List item
            self.current_list_item = {'text': '', 'level': 0}

        elif tag == 'div' and self.current_section is not None:
            # Check for special classes
            classes = attrs_dict.get('class', '').split()
            if 'placeholder' in classes:
                # Placeholder for chart/table/image
                self.current_element = {
                    'type': 'placeholder',
                    'data_type': attrs_dict.get('data-type', 'unknown'),
                    'text': ''
                }
            elif 'row' in classes:
                # Two-column layout
                self.current_element = {
                    'type': 'row',
                    'columns': []
                }
            elif 'col' in classes and self.current_element and self.current_element.get('type') == 'row':
                # Column within row
                col = {'elements': []}
                self.current_element['columns'].append(col)

    def handle_endtag(self, tag):
        if self.tag_stack:
            self.tag_stack.pop()

        if tag == 'section' and self.current_section is not None:
            # End of slide
            self.slides.append(self.current_section)
            self.current_section = None

        elif tag in ('h1', 'h2', 'h3') and self.current_element:
            # End of heading
            if self.current_element['level'] == 1:
                # H1 is the slide title
                self.current_section['title'] = self.current_element['text'].strip()
            else:
                # H2/H3 are content headings
                self.current_section['elements'].append(self.current_element)
            self.current_element = None

        elif tag == 'p' and self.current_element:
            # End of paragraph
            self.current_section['elements'].append(self.current_element)
            self.current_element = None

        elif tag in ('ul', 'ol') and self.current_list:
            # End of list
            self.current_section['elements'].append(self.current_list)
            self.current_list = None

        elif tag == 'li' and self.current_list_item:
            # End of list item
            self.current_list['items'].append(self.current_list_item)
            self.current_list_item = None

        elif tag == 'div' and self.current_element:
            # End of div (placeholder or row)
            if self.current_element.get('type') in ('placeholder', 'row'):
                self.current_section['elements'].append(self.current_element)
            self.current_element = None

    def handle_data(self, data):
        text = data.strip()
        if not text:
            return

        # Accumulate text into current element
        if self.current_list_item is not None:
            self.current_list_item['text'] += data
        elif self.current_element is not None:
            if 'text' in self.current_element:
                self.current_element['text'] += data

def _parse_html_sections(html: str) -> List[Tuple[str, Dict[str, Any]]]:
    """
    Parse HTML with <section> elements into slide structures.
    Returns: [(slide_title, slide_content_dict), ...]
    """
    parser = HTMLSlideParser()
    parser.feed(html)

    slides = []
    for section in parser.slides:
        title = section['title'] or 'Slide'
        content = {
            'elements': section['elements']
        }
        slides.append((title, content))

    return slides

def _render_html_element_to_slide(
        slide,
        element: Dict[str, Any],
        top_anchor,
        y_offset_in: float,
        avail_height_in: float,
        sources_map: Dict[int, Dict[str, str]],
        resolve_citations: bool
) -> float:
    """
    Render a single HTML element to a slide.
    Returns: new y_offset_in after placing the element.
    """
    elem_type = element.get('type')

    if elem_type == 'heading':
        # H2/H3 header
        level = element.get('level', 2)
        text = element.get('text', '').strip()

        if level == 2:
            # H2 - section header
            header_h = 0.42
            if y_offset_in + header_h > avail_height_in:
                return -1  # Signal: need new slide

            tbox = _add_textbox(
                slide,
                PAGE["content_left"],
                top_anchor + Inches(y_offset_in),
                PAGE["content_width"],
                Inches(header_h - 0.10)
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            _style_paragraph(p, space_after=Pt(4))
            r = p.add_run()
            r.text = text
            _style_run(r, size=TYPE_SCALE["slide_title"], bold=True, color=PALETTE["fg"])

            return y_offset_in + header_h + BLOCK_GAP_IN

        elif level == 3:
            # H3 - subsection header
            header_h = 0.35
            if y_offset_in + header_h > avail_height_in:
                return -1

            tbox = _add_textbox(
                slide,
                PAGE["content_left"],
                top_anchor + Inches(y_offset_in),
                PAGE["content_width"],
                Inches(header_h - 0.08)
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            _style_paragraph(p, space_after=Pt(2))
            r = p.add_run()
            r.text = text
            _style_run(r, size=TYPE_SCALE["h3"], bold=True, color=PALETTE["fg"])

            return y_offset_in + header_h + BLOCK_GAP_IN

    elif elem_type == 'paragraph':
        # Text paragraph
        text = element.get('text', '').strip()
        if not text:
            return y_offset_in

        # Estimate height
        lines = [text]  # Single paragraph
        est_h = _estimate_text_block_height(lines, PAGE["content_width"].inches, TYPE_SCALE["body"].pt)

        if y_offset_in + est_h > avail_height_in:
            return -1  # Need new slide

        tbox = _add_textbox(
            slide,
            PAGE["content_left"],
            top_anchor + Inches(y_offset_in),
            PAGE["content_width"],
            Inches(est_h)
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        p = tf.paragraphs[0]
        _style_paragraph(p, space_after=Pt(PARA_SPACE_AFTER_PT))
        _emit_text_with_formatting(p, text, sources_map, resolve_citations)

        return y_offset_in + est_h + BLOCK_GAP_IN

    elif elem_type == 'list':
        # Bullet list or numbered list
        items = element.get('items', [])
        if not items:
            return y_offset_in

        # Convert to lines format
        lines = []
        for item in items:
            text = item.get('text', '').strip()
            level = item.get('level', 0)
            if text:
                # Format as markdown-style bullet
                indent = "  " * level
                lines.append(f"{indent}- {text}")

        if not lines:
            return y_offset_in

        est_h = _estimate_text_block_height(lines, PAGE["content_width"].inches, TYPE_SCALE["body"].pt)

        if y_offset_in + est_h > avail_height_in:
            return -1  # Need new slide

        tbox = _add_textbox(
            slide,
            PAGE["content_left"],
            top_anchor + Inches(y_offset_in),
            PAGE["content_width"],
            Inches(est_h)
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        p = tf.paragraphs[0]
        first = True
        for line in lines:
            if not first:
                p = tf.add_paragraph()
            first = False
            lvl, txt = _parse_bullet_level(line)
            _style_paragraph(p, level=lvl, space_after=Pt(PARA_SPACE_AFTER_PT))
            _emit_text_with_formatting(p, txt, sources_map, resolve_citations)

        return y_offset_in + est_h + BLOCK_GAP_IN

    elif elem_type == 'placeholder':
        # Placeholder for chart/table/image
        data_type = element.get('data_type', 'unknown')
        text = element.get('text', f'{data_type.title()} Placeholder').strip()

        placeholder_h = 2.5  # Fixed height for placeholders

        if y_offset_in + placeholder_h > avail_height_in:
            return -1  # Need new slide

        # Create a rounded rectangle placeholder
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            PAGE["content_left"],
            top_anchor + Inches(y_offset_in),
            PAGE["content_width"],
            Inches(placeholder_h - 0.1)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(240, 244, 252)
        rect.line.color.rgb = PALETTE["rule"]
        rect.line.width = Pt(1.5)
        rect.line.dash_style = 2  # Dashed line

        tf = rect.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        _style_run(r, size=TYPE_SCALE["body"], italic=True, color=PALETTE["muted"])

        return y_offset_in + placeholder_h + BLOCK_GAP_IN

    # Unknown element type
    return y_offset_in

def _render_html_slide(
        prs: Presentation,
        title: str,
        content: Dict[str, Any],
        sources_map: Dict[int, Dict[str, str]],
        resolve_citations: bool
) -> None:
    """
    Render a single HTML-parsed slide.
    Handles slide splitting if content doesn't fit.
    """
    elements = content.get('elements', [])
    if not elements:
        # Empty slide with just title
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_slide_title(slide, title)
        return

    # Initial slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, title)

    top_anchor = PAGE["content_top"] + Inches(TITLE_BAND_IN)
    bottom_anchor = PAGE["slide_height"] - Inches(BOTTOM_MARGIN_IN)
    avail_height_in = _emu_to_in(bottom_anchor - top_anchor)

    y_offset_in = 0.0

    for elem in elements:
        new_y = _render_html_element_to_slide(
            slide, elem, top_anchor, y_offset_in, avail_height_in,
            sources_map, resolve_citations
        )

        if new_y == -1:
            # Element doesn't fit, start new slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_slide_title(slide, f"{title} (cont.)")
            y_offset_in = 0.0

            # Try rendering on new slide
            new_y = _render_html_element_to_slide(
                slide, elem, top_anchor, y_offset_in, avail_height_in,
                sources_map, resolve_citations
            )

            if new_y == -1:
                # Still doesn't fit - element is too large
                # Render what we can
                new_y = y_offset_in + 0.5  # Skip past it

        y_offset_in = new_y

def _normalize_html_citations(html: str) -> str:
    """
    Convert HTML citation tags to [[S:n]] format for uniform handling.

    Transforms:
        <sup class="cite" data-sids="1,3">[S:1,3]</sup> → [[S:1,3]]
        <sup class="cite">[S:1,3]</sup> → [[S:1,3]]

    Also handles footnotes-style references if present.
    """
    import re

    # Pattern 1: <sup class="cite" ...>[S:...]</sup>
    # Extract the [S:...] part and convert to [[S:...]]
    html = re.sub(
        r'<sup\s+[^>]*class=["\'][^"\']*\bcite\b[^"\']*["\'][^>]*>\s*\[S:([^\]]+)\]\s*</sup>',
        r'[[S:\1]]',
        html,
        flags=re.I
    )

    # Pattern 2: Any remaining <sup class="cite"> without the [S:...] format
    # (less common, but handle gracefully)
    html = re.sub(
        r'<sup\s+[^>]*class=["\'][^"\']*\bcite\b[^"\']*["\'][^>]*>([^<]*)</sup>',
        r'\1',  # Just strip the tags, preserve text
        html,
        flags=re.I
    )

    return html

class CSSStyleExtractor:
    """Extract and parse CSS styles from HTML <style> tags."""

    def __init__(self):
        self.styles = {}  # class_name -> {property: value}

    def parse_style_tag(self, css_text: str):
        """Parse CSS and extract class-based rules."""
        try:
            rules = tinycss2.parse_stylesheet(css_text, skip_comments=True, skip_whitespace=True)

            for rule in rules:
                if rule.type == 'qualified-rule':
                    # Extract selector (class name)
                    selector = self._get_selector(rule.prelude)
                    if not selector:
                        continue

                    # Extract declarations
                    declarations = tinycss2.parse_declaration_list(rule.content)
                    styles = {}

                    for decl in declarations:
                        if decl.type == 'declaration':
                            prop = decl.name
                            value = self._get_value(decl.value)
                            if value:
                                styles[prop] = value

                    if styles:
                        self.styles[selector] = styles
        except Exception as e:
            print(f"CSS parsing error: {e}")

    def _get_selector(self, tokens) -> Optional[str]:
        """Extract class name from selector tokens."""
        for token in tokens:
            if hasattr(token, 'value'):
                # Remove dots and spaces to get clean class name
                selector = token.value.strip().lstrip('.')
                if selector:
                    return selector
        return None

    def _get_value(self, tokens) -> Optional[str]:
        """Convert token list to string value."""
        parts = []
        for token in tokens:
            if hasattr(token, 'value'):
                parts.append(str(token.value))
            elif hasattr(token, 'serialize'):
                parts.append(token.serialize())
        return ''.join(parts).strip() if parts else None

    def get_color(self, class_name: str, property: str = 'color') -> Optional[RGBColor]:
        """Get color from class styles."""
        if class_name not in self.styles:
            return None

        value = self.styles[class_name].get(property)
        if not value:
            return None

        return self._parse_color(value)

    def get_background_color(self, class_name: str) -> Optional[RGBColor]:
        """Get background color from class styles."""
        return self.get_color(class_name, 'background')

    def _parse_color(self, color_str: str) -> Optional[RGBColor]:
        """Parse CSS color (hex or rgb) to RGBColor."""
        color_str = color_str.strip()

        # Hex color: #RRGGBB
        if color_str.startswith('#'):
            hex_color = color_str.lstrip('#')
            if len(hex_color) == 6:
                try:
                    r = int(hex_color[0:2], 16)
                    g = int(hex_color[2:4], 16)
                    b = int(hex_color[4:6], 16)
                    return RGBColor(r, g, b)
                except:
                    pass

        # rgb(r, g, b)
        rgb_match = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
        if rgb_match:
            r, g, b = map(int, rgb_match.groups())
            return RGBColor(r, g, b)

        return None


class EnhancedHTMLSlideParser(HTMLParser):
    """Parse HTML with CSS styling support."""

    def __init__(self, css_extractor: CSSStyleExtractor):
        super().__init__()
        self.css = css_extractor
        self.slides = []
        self.current_section = None
        self.current_element = None
        self.current_list = None
        self.current_list_item = None
        self.tag_stack = []

        # Formatting state stack
        self.format_stack = []  # [{bold, italic, color, bg_color, classes}]

    def _current_format(self) -> Dict[str, Any]:
        """Get current formatting state."""
        if not self.format_stack:
            return {'bold': False, 'italic': False, 'color': None, 'bg_color': None, 'classes': []}
        return self.format_stack[-1].copy()

    def _push_format(self, **kwargs):
        """Push new formatting state."""
        fmt = self._current_format()
        fmt.update(kwargs)
        self.format_stack.append(fmt)

    def _pop_format(self):
        """Pop formatting state."""
        if self.format_stack:
            self.format_stack.pop()

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        self.tag_stack.append(tag)

        # Track formatting tags
        if tag in ('strong', 'b'):
            self._push_format(bold=True)
        elif tag in ('em', 'i'):
            self._push_format(italic=True)
        elif tag == 'span':
            # Handle inline styles or classes
            style = attrs_dict.get('style', '')
            classes = attrs_dict.get('class', '').split()

            color = None
            bg_color = None

            # Parse inline color
            if 'color:' in style:
                m = re.search(r'color:\s*([^;]+)', style)
                if m:
                    color = self.css._parse_color(m.group(1))

            # Or get from class
            if classes and not color:
                for cls in classes:
                    color = self.css.get_color(cls)
                    if color:
                        break

            self._push_format(color=color, bg_color=bg_color, classes=classes)

        # Section handling
        if tag == 'section':
            self.current_section = {
                'id': attrs_dict.get('id', ''),
                'title': '',
                'elements': []
            }

        elif tag in ('h1', 'h2', 'h3') and self.current_section is not None:
            self.current_element = {
                'type': 'heading',
                'level': int(tag[1]),
                'text': '',
                'runs': []  # Track formatted runs
            }

        elif tag == 'p' and self.current_section is not None:
            classes = attrs_dict.get('class', '').split()
            self.current_element = {
                'type': 'paragraph',
                'text': '',
                'runs': [],
                'classes': classes
            }

        elif tag in ('ul', 'ol') and self.current_section is not None:
            self.current_list = {
                'type': 'list',
                'ordered': tag == 'ol',
                'items': []
            }

        elif tag == 'li' and self.current_list is not None:
            self.current_list_item = {
                'text': '',
                'runs': [],
                'level': 0
            }

        elif tag == 'div' and self.current_section is not None:
            classes = attrs_dict.get('class', '').split()

            if 'highlight' in classes:
                # Special highlighted box
                self.current_element = {
                    'type': 'highlight',
                    'text': '',
                    'runs': [],
                    'classes': classes
                }
            elif 'placeholder' in classes:
                self.current_element = {
                    'type': 'placeholder',
                    'data_type': attrs_dict.get('data-type', 'unknown'),
                    'text': ''
                }

    def handle_endtag(self, tag):
        if self.tag_stack:
            self.tag_stack.pop()

        # Pop formatting state
        if tag in ('strong', 'b', 'em', 'i', 'span'):
            self._pop_format()

        if tag == 'section' and self.current_section is not None:
            self.slides.append(self.current_section)
            self.current_section = None

        elif tag in ('h1', 'h2', 'h3') and self.current_element:
            if self.current_element['level'] == 1:
                self.current_section['title'] = self.current_element['text'].strip()
            else:
                self.current_section['elements'].append(self.current_element)
            self.current_element = None

        elif tag == 'p' and self.current_element:
            self.current_section['elements'].append(self.current_element)
            self.current_element = None

        elif tag in ('ul', 'ol') and self.current_list:
            self.current_section['elements'].append(self.current_list)
            self.current_list = None

        elif tag == 'li' and self.current_list_item:
            self.current_list['items'].append(self.current_list_item)
            self.current_list_item = None

        elif tag == 'div' and self.current_element:
            if self.current_element.get('type') in ('highlight', 'placeholder'):
                self.current_section['elements'].append(self.current_element)
            self.current_element = None

    def handle_data(self, data):
        text = data.strip()
        if not text:
            return

        # Get current formatting
        fmt = self._current_format()

        # Create a run with formatting
        run = {
            'text': data,
            'bold': fmt.get('bold', False),
            'italic': fmt.get('italic', False),
            'color': fmt.get('color'),
            'classes': fmt.get('classes', [])
        }

        # Accumulate into current element
        if self.current_list_item is not None:
            self.current_list_item['text'] += data
            self.current_list_item['runs'].append(run)
        elif self.current_element is not None:
            if 'text' in self.current_element:
                self.current_element['text'] += data
            if 'runs' in self.current_element:
                self.current_element['runs'].append(run)

# Update color parsing to handle hex without # symbol
def _parse_color(self, color_str: str) -> Optional[RGBColor]:
    """Parse CSS color (hex or rgb) to RGBColor."""
    color_str = color_str.strip()

    # Hex color with #: #RRGGBB
    if color_str.startswith('#'):
        hex_color = color_str.lstrip('#')
        if len(hex_color) == 6:
            try:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)
            except:
                pass

    # Hex color WITHOUT #: RRGGBB (CSS parser strips the #)
    elif len(color_str) == 6 and all(c in '0123456789abcdefABCDEF' for c in color_str):
        try:
            r = int(color_str[0:2], 16)
            g = int(color_str[2:4], 16)
            b = int(color_str[4:6], 16)
            return RGBColor(r, g, b)
        except:
            pass

    # 3-digit hex: FFF -> FFFFFF
    elif len(color_str) == 3 and all(c in '0123456789abcdefABCDEF' for c in color_str):
        try:
            r = int(color_str[0] * 2, 16)
            g = int(color_str[1] * 2, 16)
            b = int(color_str[2] * 2, 16)
            return RGBColor(r, g, b)
        except:
            pass

    # rgb(r, g, b)
    rgb_match = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
    if rgb_match:
        r, g, b = map(int, rgb_match.groups())
        return RGBColor(r, g, b)

    return None

def _parse_html_sections_with_css(html: str) -> Tuple[List[Tuple[str, Dict[str, Any]]], CSSStyleExtractor]:
    """Parse HTML with CSS styling support."""

    # Extract CSS
    css_extractor = CSSStyleExtractor()
    style_match = re.search(r'<style[^>]*>(.*?)</style>', html, re.DOTALL | re.IGNORECASE)
    if style_match:
        css_text = style_match.group(1)
        css_extractor.parse_style_tag(css_text)
        # 🔍 DEBUG: Print extracted styles
        print("=" * 60)
        print("EXTRACTED CSS STYLES:")
        for selector, styles in css_extractor.styles.items():
            print(f"  .{selector}:")
            for prop, value in styles.items():
                print(f"    {prop}: {value}")
        print("=" * 60)

    # Parse HTML
    parser = EnhancedHTMLSlideParser(css_extractor)
    parser.feed(html)

    slides = []
    for section in parser.slides:
        title = section['title'] or 'Slide'
        content = {
            'elements': section['elements']
        }
        slides.append((title, content))

    return slides, css_extractor


def _emit_formatted_runs(paragraph, runs: List[Dict[str, Any]], sources_map: Dict[int, Dict[str, str]], resolve_citations: bool, css: CSSStyleExtractor):
    """Emit text runs with formatting - matching HTML visual design."""

    for run_data in runs:
        text = run_data.get('text', '')
        if not text:
            continue

        # Get formatting from this run
        bold = run_data.get('bold', False)
        italic = run_data.get('italic', False)
        color = run_data.get('color')

        # CRITICAL: In the HTML, <strong> tags are BLUE (#0066cc)
        # So if text is bold, make it blue (like the original HTML design)
        if bold and not color:
            color = css.get_color('strong') or RGBColor(0, 102, 204)  # Blue for bold

        # Get color from classes if still not set
        if not color:
            for cls in run_data.get('classes', []):
                color = css.get_color(cls)
                if color:
                    break

        # Default color (dark gray from HTML body text)
        if not color:
            color = RGBColor(51, 51, 51)  # #333 from HTML

        # Handle citations while preserving formatting
        if resolve_citations and '[[S:' in text:
            parts = re.split(r'(\[\[S:\d+(?:,\d+)*\]\])', text)
            for part in parts:
                if not part:
                    continue

                cit_match = _CIT_RE.match(part)
                if cit_match:
                    sid = int(cit_match.group(1))
                    rec = sources_map.get(sid, {})
                    label = rec.get("title") or f"[{sid}]"
                    url = rec.get("url", "")

                    r = paragraph.add_run()
                    r.text = label
                    _style_run(r, size=TYPE_SCALE["body"], bold=bold, italic=italic, color=PALETTE["accent"])
                    if url:
                        try:
                            r.hyperlink.address = url
                        except:
                            pass
                else:
                    r = paragraph.add_run()
                    r.text = part
                    _style_run(r, size=TYPE_SCALE["body"], bold=bold, italic=italic, color=color)
        else:
            r = paragraph.add_run()
            r.text = text
            _style_run(r, size=TYPE_SCALE["body"], bold=bold, italic=italic, color=color)

def _render_html_element_to_slide_styled(
        slide,
        element: Dict[str, Any],
        top_anchor,
        y_offset_in: float,
        avail_height_in: float,
        sources_map: Dict[int, Dict[str, str]],
        resolve_citations: bool,
        css: CSSStyleExtractor
) -> float:
    """Render HTML element with CSS styling."""

    elem_type = element.get('type')

    if elem_type == 'heading':
        level = element.get('level', 2)
        runs = element.get('runs', [])

        if level == 2:
            # H2 - section header - MAKE IT BLUE
            header_h = 0.42
            if y_offset_in + header_h > avail_height_in:
                return -1

            tbox = _add_textbox(
                slide,
                PAGE["content_left"],
                top_anchor + Inches(y_offset_in),
                PAGE["content_width"],
                Inches(header_h - 0.10)
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            _style_paragraph(p, space_after=Pt(4))

            # Emit formatted runs with COLOR from CSS
            if runs:
                for run_data in runs:
                    text = run_data.get('text', '').strip()
                    if not text:
                        continue
                    r = p.add_run()
                    r.text = text
                    # Get H2 color from CSS (should be #0066cc)
                    h2_color = css.get_color('h2') or RGBColor(0, 102, 204)  # Default blue
                    _style_run(r, size=TYPE_SCALE["slide_title"], bold=True, color=h2_color)
            else:
                # Fallback
                text = element.get('text', '').strip()
                r = p.add_run()
                r.text = text
                h2_color = css.get_color('h2') or RGBColor(0, 102, 204)
                _style_run(r, size=TYPE_SCALE["slide_title"], bold=True, color=h2_color)

            return y_offset_in + header_h + BLOCK_GAP_IN

        elif level == 3:
            header_h = 0.35
            if y_offset_in + header_h > avail_height_in:
                return -1

            tbox = _add_textbox(
                slide,
                PAGE["content_left"],
                top_anchor + Inches(y_offset_in),
                PAGE["content_width"],
                Inches(header_h - 0.08)
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            _style_paragraph(p, space_after=Pt(2))

            _emit_formatted_runs(p, runs, sources_map, resolve_citations, css)

            return y_offset_in + header_h + BLOCK_GAP_IN

    elif elem_type == 'paragraph':
        runs = element.get('runs', [])
        classes = element.get('classes', [])

        if not runs:
            return y_offset_in

        # Estimate height
        text = element.get('text', '')
        lines = [text]
        est_h = _estimate_text_block_height(lines, PAGE["content_width"].inches, TYPE_SCALE["body"].pt)

        if y_offset_in + est_h > avail_height_in:
            return -1

        tbox = _add_textbox(
            slide,
            PAGE["content_left"],
            top_anchor + Inches(y_offset_in),
            PAGE["content_width"],
            Inches(est_h)
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        p = tf.paragraphs[0]
        _style_paragraph(p, space_after=Pt(PARA_SPACE_AFTER_PT))

        # Emit formatted runs
        _emit_formatted_runs(p, runs, sources_map, resolve_citations, css)

        return y_offset_in + est_h + BLOCK_GAP_IN

    elif elem_type == 'highlight':
        # Special highlighted box with background color
        runs = element.get('runs', [])
        classes = element.get('classes', [])

        if not runs:
            return y_offset_in

        text = element.get('text', '')
        lines = [text]
        est_h = _estimate_text_block_height(lines, PAGE["content_width"].inches, TYPE_SCALE["body"].pt)
        est_h = max(0.6, est_h + 0.25)  # Add padding

        if y_offset_in + est_h > avail_height_in:
            return -1

        # Get background color from CSS
        bg_color = RGBColor(230, 242, 255)  # Default light blue
        for cls in classes:
            parsed_bg = css.get_background_color(cls)
            if parsed_bg:
                bg_color = parsed_bg
                break

        # Create colored box
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            PAGE["content_left"],
            top_anchor + Inches(y_offset_in),
            PAGE["content_width"],
            Inches(est_h - 0.08)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg_color

        # Border
        border_color = css.get_color('highlight', 'border-left-color') or PALETTE["accent"]
        rect.line.color.rgb = border_color
        rect.line.width = Pt(3)

        tf = rect.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.12)

        p = tf.paragraphs[0]
        _style_paragraph(p, space_after=Pt(2))

        # Emit formatted runs
        _emit_formatted_runs(p, runs, sources_map, resolve_citations, css)

        return y_offset_in + est_h + BLOCK_GAP_IN

    elif elem_type == 'list':
        items = element.get('items', [])
        if not items:
            return y_offset_in

        # Estimate total height
        total_text = '\n'.join(item.get('text', '') for item in items)
        lines = [f"- {item.get('text', '')}" for item in items]
        est_h = _estimate_text_block_height(lines, PAGE["content_width"].inches, TYPE_SCALE["body"].pt)

        if y_offset_in + est_h > avail_height_in:
            return -1

        tbox = _add_textbox(
            slide,
            PAGE["content_left"],
            top_anchor + Inches(y_offset_in),
            PAGE["content_width"],
            Inches(est_h)
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        p = tf.paragraphs[0]
        first = True

        for item in items:
            if not first:
                p = tf.add_paragraph()
            first = False

            level = item.get('level', 0)
            runs = item.get('runs', [])

            _style_paragraph(p, level=level, space_after=Pt(PARA_SPACE_AFTER_PT))
            _emit_formatted_runs(p, runs, sources_map, resolve_citations, css)

        return y_offset_in + est_h + BLOCK_GAP_IN

    # Fallback
    return y_offset_in


def _render_html_slide_styled(
        prs: Presentation,
        title: str,
        content: Dict[str, Any],
        sources_map: Dict[int, Dict[str, str]],
        resolve_citations: bool,
        css: CSSStyleExtractor
) -> None:
    """Render HTML slide with CSS styling."""

    elements = content.get('elements', [])
    if not elements:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_slide_title(slide, title)
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, title)

    top_anchor = PAGE["content_top"] + Inches(TITLE_BAND_IN)
    bottom_anchor = PAGE["slide_height"] - Inches(BOTTOM_MARGIN_IN)
    avail_height_in = _emu_to_in(bottom_anchor - top_anchor)

    y_offset_in = 0.0

    for elem in elements:
        new_y = _render_html_element_to_slide_styled(
            slide, elem, top_anchor, y_offset_in, avail_height_in,
            sources_map, resolve_citations, css
        )

        if new_y == -1:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_slide_title(slide, f"{title} (cont.)")
            y_offset_in = 0.0

            new_y = _render_html_element_to_slide_styled(
                slide, elem, top_anchor, y_offset_in, avail_height_in,
                sources_map, resolve_citations, css
            )

            if new_y == -1:
                new_y = y_offset_in + 0.5

        y_offset_in = new_y

## #
def _create_html_styled_slide(prs: Presentation, title: str) -> tuple:
    """
    Create a slide that matches the HTML visual design:
    - Gray background (#f5f5f5)
    - White content box with shadow
    - Blue underline under title
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Set slide background to light gray (like HTML body)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)  # #f5f5f5

    # Create white content box (like HTML <section>)
    content_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4),  # Left margin
        Inches(0.4),  # Top margin
        Inches(9.2),  # Width (10" - 0.4" margins)
        Inches(6.7),  # Height
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

    # Subtle shadow effect (approximate)
    content_box.line.color.rgb = RGBColor(220, 220, 220)
    content_box.line.width = Pt(0.5)

    # Add title with blue underline
    title_top = Inches(0.7)
    title_left = Inches(0.8)
    title_width = Inches(8.4)

    # Title text
    title_box = slide.shapes.add_textbox(
        title_left,
        title_top,
        title_width,
        Inches(0.8)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.LEFT

    title_run = title_p.add_run()
    title_run.text = title
    title_run.font.size = Pt(32)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(26, 26, 26)  # #1a1a1a

    # Blue underline bar (like HTML h1 border-bottom)
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        title_left,
        title_top + Inches(0.75),
        title_width,
        Pt(4)  # 4px thick
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = RGBColor(0, 102, 204)  # #0066cc
    underline.line.fill.background()

    # Calculate available space for content
    top_anchor = title_top + Inches(1.0)
    bottom_anchor = Inches(6.7)
    available_height_in = _emu_to_in(bottom_anchor - top_anchor)

    return slide, top_anchor, available_height_in, 0.0, title_left, title_width

def _render_html_element_styled(
        slide,
        element: Dict[str, Any],
        top_anchor,
        y_offset_in: float,
        avail_height_in: float,
        sources_map: Dict[int, Dict[str, str]],
        resolve_citations: bool,
        css: CSSStyleExtractor,
        content_left,
        content_width
) -> float:
    """Render HTML element with full visual styling."""

    elem_type = element.get('type')

    if elem_type == 'heading':
        level = element.get('level', 2)
        runs = element.get('runs', [])

        if level == 2:
            # H2 - Blue heading
            header_h = 0.42
            if y_offset_in + header_h > avail_height_in:
                return -1

            tbox = slide.shapes.add_textbox(
                content_left,
                top_anchor + Inches(y_offset_in),
                content_width,
                Inches(header_h - 0.10)
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            p.space_after = Pt(6)
            p.alignment = PP_ALIGN.LEFT

            # H2 is BLUE in HTML
            h2_color = css.get_color('h2') or RGBColor(0, 102, 204)

            if runs:
                for run_data in runs:
                    text = run_data.get('text', '').strip()
                    if text:
                        r = p.add_run()
                        r.text = text
                        r.font.size = Pt(22)
                        r.font.bold = True
                        r.font.color.rgb = h2_color

            return y_offset_in + header_h + BLOCK_GAP_IN

    elif elem_type == 'paragraph':
        runs = element.get('runs', [])
        if not runs:
            return y_offset_in

        text = element.get('text', '')
        lines = [text]
        est_h = _estimate_text_block_height(lines, content_width.inches, TYPE_SCALE["body"].pt)

        if y_offset_in + est_h > avail_height_in:
            return -1

        tbox = slide.shapes.add_textbox(
            content_left,
            top_anchor + Inches(y_offset_in),
            content_width,
            Inches(est_h)
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        p = tf.paragraphs[0]
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT

        _emit_formatted_runs(p, runs, sources_map, resolve_citations, css)

        return y_offset_in + est_h + BLOCK_GAP_IN

    elif elem_type == 'highlight':
        # Highlight box with light blue background
        runs = element.get('runs', [])
        if not runs:
            return y_offset_in

        text = element.get('text', '')
        lines = [text]
        est_h = _estimate_text_block_height(lines, content_width.inches, TYPE_SCALE["body"].pt)
        est_h = max(0.6, est_h + 0.3)

        if y_offset_in + est_h > avail_height_in:
            return -1

        # Light blue background (#e6f2ff)
        bg_color = css.get_background_color('highlight') or RGBColor(230, 242, 255)

        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            content_left,
            top_anchor + Inches(y_offset_in),
            content_width,
            Inches(est_h - 0.08)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg_color
        rect.line.fill.background()

        # Blue left border (4px solid #0066cc)
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            content_left,
            top_anchor + Inches(y_offset_in),
            Pt(4),
            Inches(est_h - 0.08)
        )
        border.fill.solid()
        border.fill.fore_color.rgb = RGBColor(0, 102, 204)
        border.line.fill.background()

        tf = rect.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.12)

        p = tf.paragraphs[0]
        p.space_after = Pt(2)

        _emit_formatted_runs(p, runs, sources_map, resolve_citations, css)

        return y_offset_in + est_h + BLOCK_GAP_IN

    elif elem_type == 'list':
        items = element.get('items', [])
        if not items:
            return y_offset_in

        lines = [f"- {item.get('text', '')}" for item in items]
        est_h = _estimate_text_block_height(lines, content_width.inches, TYPE_SCALE["body"].pt)

        if y_offset_in + est_h > avail_height_in:
            return -1

        tbox = slide.shapes.add_textbox(
            content_left,
            top_anchor + Inches(y_offset_in),
            content_width,
            Inches(est_h)
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        p = tf.paragraphs[0]
        first = True

        for item in items:
            if not first:
                p = tf.add_paragraph()
            first = False

            level = item.get('level', 0)
            runs = item.get('runs', [])

            p.level = level
            p.space_after = Pt(6)
            p.alignment = PP_ALIGN.LEFT

            _emit_formatted_runs(p, runs, sources_map, resolve_citations, css)

        return y_offset_in + est_h + BLOCK_GAP_IN

    return y_offset_in

def _render_html_slide_styled(
        prs: Presentation,
        title: str,
        content: Dict[str, Any],
        sources_map: Dict[int, Dict[str, str]],
        resolve_citations: bool,
        css: CSSStyleExtractor
) -> None:
    """Render HTML slide matching the original HTML visual design."""

    elements = content.get('elements', [])
    if not elements:
        slide, top_anchor, avail_height_in, y_offset_in, left, width = _create_html_styled_slide(prs, title)
        return

    slide, top_anchor, avail_height_in, y_offset_in, content_left, content_width = _create_html_styled_slide(prs, title)

    for elem in elements:
        new_y = _render_html_element_styled(
            slide, elem, top_anchor, y_offset_in, avail_height_in,
            sources_map, resolve_citations, css, content_left, content_width
        )

        if new_y == -1:
            # Need new slide
            slide, top_anchor, avail_height_in, y_offset_in, content_left, content_width = _create_html_styled_slide(
                prs, f"{title} (cont.)"
            )

            new_y = _render_html_element_styled(
                slide, elem, top_anchor, y_offset_in, avail_height_in,
                sources_map, resolve_citations, css, content_left, content_width
            )

            if new_y == -1:
                new_y = y_offset_in + 0.5

        y_offset_in = new_y
##

def render_pptx(
        path: str,
        content_md: str = "",
        content_html: str = "",
        *,
        title: Optional[str] = None,
        base_dir: Optional[str] = None,
        sources: Optional[str] = None,
        resolve_citations: bool = False,
        include_sources_slide: bool = False
) -> str:
    """
    Render PPTX from Markdown OR HTML.

    HTML mode (preferred):
    - Expects <section> elements, one per slide
    - <h1> or <h2> = slide title
    - <p>, <ul>, <div class="placeholder"> = content
    - Supports [[S:n]] citations in text

    Markdown mode (fallback):
    - Uses ## per-slide headings
    - Standard markdown syntax

    Args:
        content_md: Markdown content (legacy)
        content_html: HTML content (preferred)
        title: Optional presentation title (for title slide)
        sources: JSON sources for citations
        resolve_citations: Convert [[S:n]] to hyperlinks
        include_sources_slide: Add sources slide at end

    Returns:
        Saved PPTX filename (basename only)
    """
    basename = _basename_only(path, ".pptx")
    outdir = _outdir()
    outfile = outdir / basename
    _ensure_parent(outfile)

    sources_map: Dict[int, Dict[str, str]] = {}
    order: List[int] = []
    if sources:
        sources_map, order = md_utils._normalize_sources(sources)

    prs = Presentation()

    # Determine rendering mode
    use_html = bool(content_html and content_html.strip())

    if use_html:
        # Parse HTML with CSS support
        normalized_html = _normalize_html_citations(content_html)
        sections, css_extractor = _parse_html_sections_with_css(normalized_html)

        # Title slide
        if title:
            _add_title_slide(prs, title)
        elif sections:
            # Use first slide title
            first_title, _ = sections[0]
            _add_title_slide(prs, first_title)

        # Content slides with styling
        for slide_title, slide_content in sections:
            _render_html_slide_styled(prs, slide_title, slide_content, sources_map, resolve_citations, css_extractor)

    else:
        # Markdown mode (existing logic)
        sections = _split_markdown_sections(content_md or "")

        # Title slide
        if title:
            _add_title_slide(prs, title)
        else:
            fst_title, _ = sections[0]
            _add_title_slide(prs, fst_title)

        # Content slides
        for stitle, body in sections:
            _render_section_across_slides(prs, stitle, body, sources_map, resolve_citations)

    # Sources slide
    if include_sources_slide and sources_map:
        _add_sources_slide(prs, sources_map, order)

    prs.save(str(outfile))
    return basename